import os
from pptx import Presentation

def pptx_to_markdown(pptx_path, output_folder):
    # Ensure output folder exists
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Load the presentation
    prs = Presentation(pptx_path)
    
    # Placeholder for markdown content
    md_content = []
    
    # Iterate through slides and extract content
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            # Check if shape has text
            if shape.has_text_frame:
                text = shape.text_frame.text
                if shape == slide.shapes[0]:  # Assume first shape is title
                    md_content.append(f"# {text}\n")
                else:
                    md_content.append(f"## {text}\n")
            
            # Check if shape has image
            if shape.shape_type == 13:  # Type 13 corresponds to Picture
                img = shape.image
                img_bytes = img.blob
                img_format = img.ext

                # Save the image
                img_name = f"slide_{slide_number}_img.{img_format}"
                img_path = os.path.join(output_folder, img_name)
                with open(img_path, "wb") as img_file:
                    img_file.write(img_bytes)
                
                # Add reference to markdown content
                md_content.append(f"![{img_name}]({img_path})\n")
                
    # Write markdown content to a file
    md_path = os.path.join(output_folder, "output.md")
    with open(md_path, "w") as md_file:
        md_file.write("\n".join(md_content))

    print(f"Markdown saved to: {md_path}")

# Example usage
pptx_path = 'Topic 7 Cost Accounting.pptx'
output_folder = '/topic_6_img'
pptx_to_markdown(pptx_path, output_folder)
